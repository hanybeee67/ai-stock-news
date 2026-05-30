from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# 슬라이드 1: 타이틀
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AI 증시 브리핑 앱 유저 매뉴얼"
subtitle.text = "매일 아침 딱 3분, 글로벌 뉴스로 단타 전략 세우기\n(버전: v2.1.0)"

# 슬라이드 2: 핵심 기능 요약
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "1. 핵심 기능 요약 🚀"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "• 완전 자동화 분석: 매일 새벽 5시 30분, 글로벌 경제 뉴스 수집 및 AI 분석 완료"
p = tf.add_paragraph()
p.text = "• 핵심 선별 필터링: 14개 국내외 피드에서 가중치 기반 중요 뉴스 5건 선별"
p = tf.add_paragraph()
p.text = "• 나비효과 분석: 단순 요약이 아닌 1차 → 2차 → 3차 파급 효과 심층 추론"
p = tf.add_paragraph()
p.text = "• 수혜주 및 리스크 방어벽: 직관적인 수혜주 매핑과 확실한 손절 기준 제공"

# 슬라이드 3: 메인 화면 (오늘 브리핑)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "2. 앱 메인 화면 (오늘 브리핑) 📊"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "• 오늘의 AI 핵심: 그날의 전체 시장 분위기(Market Mood)와 핵심 요약 1줄 제공"
p = tf.add_paragraph()
p.text = "• 상세 뉴스 카드 탭 구성:"
p = tf.add_paragraph()
p.text = "  - [🔍 AI 분석] 찰리(월스트리트 프롭트레이더 페르소나)의 심층 분석"
p.level = 1
p = tf.add_paragraph()
p.text = "  - [📈 수혜주] 분석 근거와 함께 관련 주식(KRX/US) 및 최근 추세선 매핑"
p.level = 1
p = tf.add_paragraph()
p.text = "  - [⚠️ 리스크] 언제 손절하고 도망가야 하는지 명확히 짚어주는 단타 경고"
p.level = 1

# 슬라이드 4: 히스토리 & 저장 화면
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "3. 히스토리 & 뉴스 저장 기능 📅"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "• [📅 히스토리 탭]"
p = tf.add_paragraph()
p.text = "  - 최근 7일치 리포트를 캘린더 형태로 다시 보기"
p.level = 1
p = tf.add_paragraph()
p.text = "  - 당일 시장 상승/하락 지표 및 AI 신뢰도 요약 카드 제공"
p.level = 1
p = tf.add_paragraph()
p.text = "• [🔖 저장됨 탭]"
p = tf.add_paragraph()
p.text = "  - 뉴스 상세 화면 우측 상단의 📌 핀 아이콘으로 중요 뉴스 스크랩"
p.level = 1
p = tf.add_paragraph()
p.text = "  - 카테고리별 필터링 기능 지원 및 카드를 길게(Long Press) 눌러 삭제 가능"
p.level = 1

# 슬라이드 5: 개인화 설정
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "4. 활용 팁 & 개인 맞춤 설정 ⚙️"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "• 맞춤 알림 시간: 내 출근 시간, 기상 시간에 맞춰 푸시 알림 설정 (예: 07:30)"
p = tf.add_paragraph()
p.text = "• 관심 카테고리 필터: 반도체, 바이오, AI 등 원하는 테마 뉴스만 홈 화면에 필터링"
p = tf.add_paragraph()
p.text = "• ⚡ 수동 분석 트리거: 장중에 중요한 뉴스가 터지면 언제든 즉시 분석 수동 실행"
p = tf.add_paragraph()
p.text = "• 서버 상태 실시간 확인: 앱 내에서 AI 분석 진행률(📡수집중 → 🤖분석중) 파악 가능"

prs.save("c:/stock management/AI_Stock_Briefing_Manual.pptx")
print("✅ 매뉴얼 생성 완료!")
